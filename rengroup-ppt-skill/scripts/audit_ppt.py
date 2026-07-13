"""Audit a PPTX against the Ren Group presentation style."""

import argparse
import re
import sys
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN


DEFAULT_TEMPLATE = Path(__file__).resolve().parents[1] / "assets" / "任组PPT模板.pptx"
BLACK = "000000"
WHITE = "FFFFFF"
EMU_PER_INCH = 914400
DATE_RE = re.compile(r"(?:20\d{2}[./年-]\d{1,2}|日期|研究日期|汇报时间)")


def parse_args():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE,
                        help="Template used to derive title and conclusion-frame colors.")
    return parser.parse_args()


def template_colors(template_path):
    with ZipFile(template_path) as archive:
        master = archive.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
        title_match = re.search(
            r"<p:titleStyle>.*?<a:solidFill>\s*<a:srgbClr val=\"([0-9A-Fa-f]{6})\"",
            master,
            re.DOTALL,
        )
        title_color = title_match.group(1).upper() if title_match else "800000"

        conclusion_color = None
        for name in archive.namelist():
            if not re.fullmatch(r"ppt/slides/slide\d+\.xml", name):
                continue
            slide_xml = archive.read(name).decode("utf-8")
            frame_match = re.search(
                r"<a:prstGeom prst=\"roundRect\".*?<a:ln[^>]*>.*?"
                r"<a:solidFill>\s*<a:srgbClr val=\"([0-9A-Fa-f]{6})\"",
                slide_xml,
                re.DOTALL,
            )
            if frame_match:
                conclusion_color = frame_match.group(1).upper()
                break
    return title_color, conclusion_color or "C00000"


def explicit_rgb(color_format):
    try:
        value = color_format.rgb
        return str(value).upper() if value is not None else None
    except (AttributeError, TypeError):
        return None


def explicit_fill_rgb(fill_format):
    try:
        return explicit_rgb(fill_format.fore_color)
    except (AttributeError, TypeError):
        return None


def is_title(shape):
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        return shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except ValueError:
        return False


def iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                yield paragraph, run


def within(inner, outer):
    return (
        inner.left >= outer.left
        and inner.top >= outer.top
        and inner.left + inner.width <= outer.left + outer.width
        and inner.top + inner.height <= outer.top + outer.height
    )


def audit(path, title_color, conclusion_color):
    prs = Presentation(path)
    errors, warnings = [], []
    sw, sh = prs.slide_width, prs.slide_height

    for slide_no, slide in enumerate(prs.slides, 1):
        titles = [shape for shape in slide.shapes if is_title(shape)]
        if slide_no > 1:
            if not titles:
                errors.append(f"Slide {slide_no}: no title placeholder found.")
            for title in titles:
                center_offset = abs((title.left + title.width / 2) - sw / 2) / EMU_PER_INCH
                if center_offset > 0.25 or title.top > 1.15 * EMU_PER_INCH:
                    errors.append(f"Slide {slide_no}: title is not centered at the top.")
                for paragraph in title.text_frame.paragraphs:
                    if paragraph.text.strip() and paragraph.alignment != PP_ALIGN.CENTER:
                        warnings.append(f"Slide {slide_no}: title alignment is not explicitly centered.")
                    for run in paragraph.runs:
                        if not run.text.strip():
                            continue
                        color = explicit_rgb(run.font.color)
                        if color is None:
                            warnings.append(f"Slide {slide_no}: title color is inherited; confirm the template resolves it to #{title_color}.")
                        elif color != title_color:
                            errors.append(f"Slide {slide_no}: title color is #{color}, expected template color #{title_color}.")

        red_frames = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                line_color = explicit_rgb(shape.line.color)
                if line_color == conclusion_color:
                    red_frames.append(shape)
                    low_enough = shape.top >= sh * 0.62 or (
                        shape.left >= sw * 0.48 and shape.top >= sh * 0.45
                    )
                    if not low_enough:
                        errors.append(f"Slide {slide_no}: conclusion frame is not at the bottom or lower-right.")
                    if shape.fill.type is not None:
                        fill_color = explicit_fill_rgb(shape.fill)
                        if fill_color not in (None, WHITE):
                            errors.append(f"Slide {slide_no}: conclusion frame must have no fill.")

        conclusion_text = [
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "结论" in shape.text
        ]
        if conclusion_text and not red_frames:
            errors.append(f"Slide {slide_no}: conclusion text has no template-color #{conclusion_color} conclusion frame.")
        for text_shape in conclusion_text:
            containing = [frame for frame in red_frames if within(text_shape, frame)]
            if not containing:
                warnings.append(f"Slide {slide_no}: conclusion text is not fully inside the conclusion frame.")
            for _, run in iter_runs(text_shape):
                color = explicit_rgb(run.font.color)
                if color not in (None, BLACK):
                    errors.append(f"Slide {slide_no}: conclusion text must be black.")
                if run.font.size and not 20 <= run.font.size.pt <= 24:
                    errors.append(f"Slide {slide_no}: conclusion text is {run.font.size.pt:g} pt; expected 20–24 pt.")

        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if slide_no > 1 and text and DATE_RE.search(text):
                if shape.left < sw * 0.45 and shape.top > sh * 0.72:
                    errors.append(f"Slide {slide_no}: lower-left date/time marker is not allowed.")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    if shape.image.ext.lower() not in ("svg", "emf", "wmf"):
                        px_w, px_h = shape.image.size
                        ppi = min(px_w / (shape.width / EMU_PER_INCH), px_h / (shape.height / EMU_PER_INCH))
                        if ppi < 150:
                            warnings.append(f"Slide {slide_no}: image '{shape.name}' is only {ppi:.0f} ppi at display size.")
                except (AttributeError, ZeroDivisionError):
                    pass

            if getattr(shape, "has_table", False) and slide_no > 1:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fill_color = explicit_fill_rgb(cell.fill) if cell.fill.type is not None else None
                        if fill_color not in (None, WHITE):
                            errors.append(f"Slide {slide_no}: table cells must not use colored fills.")
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip() and run.font.size and not 14 <= run.font.size.pt <= 16:
                                    errors.append(f"Slide {slide_no}: table text is {run.font.size.pt:g} pt; expected 14–16 pt.")

            if is_title(shape) or slide_no == 1:
                continue
            for _, run in iter_runs(shape):
                size = run.font.size.pt if run.font.size else None
                if size is not None and size < 12:
                    errors.append(f"Slide {slide_no}: text '{run.text[:24]}' is below 12 pt.")
                elif size is not None and 16 < size < 20:
                    warnings.append(f"Slide {slide_no}: text '{run.text[:24]}' is {size:g} pt; use 14–16 or 20–24 pt by role.")
                elif size is not None and size > 24:
                    warnings.append(f"Slide {slide_no}: non-title text '{run.text[:24]}' exceeds 24 pt.")

                has_cn = any(ord(char) > 127 for char in run.text)
                has_en = bool(re.search(r"[A-Za-z]", run.text))
                font = run.font.name
                if font is None:
                    warnings.append(f"Slide {slide_no}: font for '{run.text[:24]}' is inherited; set it explicitly.")
                else:
                    if has_cn and font not in ("黑体", "SimHei"):
                        errors.append(f"Slide {slide_no}: Chinese text uses '{font}', expected 黑体.")
                    if has_en and font != "Arial":
                        errors.append(f"Slide {slide_no}: English text uses '{font}', expected Arial.")

            if getattr(shape, "has_chart", False):
                chart = shape.chart
                for axis_name in ("category_axis", "value_axis"):
                    try:
                        axis = getattr(chart, axis_name)
                        size = axis.tick_labels.font.size
                        if size and size.pt < 12:
                            errors.append(f"Slide {slide_no}: {axis_name} tick labels are below 12 pt.")
                    except (AttributeError, ValueError):
                        pass
                try:
                    size = chart.legend.font.size
                    if size and not 14 <= size.pt <= 16:
                        errors.append(f"Slide {slide_no}: chart legend is {size.pt:g} pt; expected 14–16 pt.")
                except (AttributeError, ValueError):
                    pass

    return errors, warnings


def main():
    args = parse_args()
    title_color, conclusion_color = template_colors(args.template)
    print(f"Template colors: title #{title_color}, conclusion frame #{conclusion_color}")
    errors, warnings = audit(args.pptx, title_color, conclusion_color)
    for message in warnings:
        print("WARN:", message)
    for message in errors:
        print("ERROR:", message)
    print(f"Audit complete: {len(errors)} error(s), {len(warnings)} warning(s).")
    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
